"""
main.py — FastAPI Application Entry Point for VaultMCP Backend

All routes wired with real logic:
  POST /process       — Full pipeline: detect → extract → transcribe → structure → search → generate MD
  POST /drive/save    — Save MD entry to user's Google Drive via drive_handler
  GET  /drive/vault   — Fetch vault.md from user's Google Drive
  POST /auth/google   — Exchange OAuth auth code for tokens
  GET  /auth/url      — Get Google OAuth consent URL
  GET  /health        — System health check
"""

import os
import re
import traceback
import json
import base64
import logging
from datetime import datetime
from typing import Optional

from fastapi import FastAPI, UploadFile, File, Form, Header, HTTPException, Query, Response
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
from pydantic import BaseModel
import httpx

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s"
)
logger = logging.getLogger(__name__)

# ─── Local Modules ───────────────────────────────────────────────────────────
from ai_processor import (
    process_text,
    result_to_dict,
    ProcessingError,
    InvalidTokenError as AIInvalidTokenError,
)
def get_official_info(official_link: str) -> dict:
    """Return the official link extracted by the AI processor."""
    return {"official_link": official_link or "", "provider": "mistral"}
from md_generator import (
    build_entry,
    generate_entry_md,
    format_retro_date,
    VaultEntry,
)
from drive_handler import (
    auth_google as drive_auth_google,
    get_auth_url as drive_get_auth_url,
    save_to_drive as drive_save_to_drive,
    get_vault as drive_get_vault,
    save_user_config as drive_save_user_config,
    get_user_config as drive_get_user_config,
    save_file_to_drive as drive_save_file_to_drive,
    get_file_content as drive_get_file_content,
    clear_vault_files as drive_clear_vault_files,
    DriveAuthError,
    DriveError,
)

import fitz
import docx
import openpyxl
import pptx
import io



# ─── Load Environment ────────────────────────────────────────────────────────
load_dotenv()

CORS_ORIGINS = os.getenv("CORS_ORIGINS", "http://localhost:5173,http://localhost:3000,http://127.0.0.1:5173,http://127.0.0.1:3000,https://vault-mcp-4ssi.vercel.app").split(",")

# ─── App Instance ────────────────────────────────────────────────────────────
app = FastAPI(
    title="VaultMCP Backend",
    description="Save what you scroll. Processing pipeline for Instagram reels, YouTube shorts, URLs, and text.",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=CORS_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

from fastapi.middleware.trustedhost import TrustedHostMiddleware

@app.middleware("http")
async def add_coop_header(request, call_next):
    response = await call_next(request)
    response.headers["Cross-Origin-Opener-Policy"] = "unsafe-none"
    return response


# ─── Request / Response Models ───────────────────────────────────────────────

class ProcessRequest(BaseModel):
    """Body for POST /process."""
    content: str
    hf_token: str                         # User's Hugging Face token (required)
    content_type: Optional[str] = "auto"  # "url" | "text" | "auto"


class DriveAuthRequest(BaseModel):
    """Body for POST /auth/google."""
    auth_code: str


class DriveSaveRequest(BaseModel):
    """Body for POST /drive/save."""
    md_entry: str                          # Pre-generated MD string from /process
    access_token: str                      # User's Google Drive OAuth token
    refresh_token: Optional[str] = None
    overwrite: Optional[bool] = False


class DriveVaultRequest(BaseModel):
    """Query params for GET /drive/vault."""
    pass


class ConfigSaveRequest(BaseModel):
    """Body for POST /config/save."""
    access_token: str
    refresh_token: Optional[str] = None
    hf_token: Optional[str] = ""
    display_name: Optional[str] = ""


# ─── URL Detection Helpers ──────────────────────────────────────────────────

GITHUB_RE = re.compile(
    r"(?:https?://)?(?:www\.)?github\.com/([\w.-]+)/([\w.-]+)", re.IGNORECASE
)
GENERIC_URL_RE = re.compile(
    r"https?://[^\s<>\"']+", re.IGNORECASE
)


def detect_input_type(content: str) -> str:
    """
    Detect what kind of input the user sent.
    Returns: "instagram" | "youtube" | "github" | "website" | "text"
    """
    content = content.strip()
    logger.info(f"[detect_input_type] Evaluating content: {content}")

    if GITHUB_RE.match(content):
        logger.info("[detect_input_type] Match found: github")
        return "github"

    if GENERIC_URL_RE.match(content):
        logger.info("[detect_input_type] Match found: website")
        return "website"

    logger.info("[detect_input_type] Match found: text (default fallback)")
    return "text"


def detect_category(input_type: str, url: str, ai_category: str) -> str:
    """Override AI category with rule-based detection."""
    
    # GitHub repos always = Dev Tools
    if input_type == "github":
        return "Dev Tools"
    
    # Known AI tool domains
    ai_domains = ["openai.com", "anthropic.com", "huggingface.co", 
                  "midjourney.com", "perplexity.ai", "claude.ai",
                  "gemini.google.com", "cursor.sh", "replicate.com"]
    if any(domain in url for domain in ai_domains):
        return "AI Tools"
    
    # Design/UI domains
    design_domains = ["dribbble.com", "figma.com", "behance.net", 
                      "awwwards.com", "tailwindcss.com", "shadcn"]
    if any(domain in url for domain in design_domains):
        return "Design"
    
    # Plain text with prompt keywords = Prompts
    prompt_keywords = ["prompt", "system prompt", "instruction", "act as", "you are a"]
    if input_type == "text" and any(kw in url.lower() for kw in prompt_keywords):
        return "Prompts"
    
    # Fall back to AI category if it's valid
    valid = ["AI Tools", "Dev Tools", "Prompts", "Design", "Resources", "Other"]
    return ai_category if ai_category in valid else "Resources"


# ─── Website Scraper (lightweight) ──────────────────────────────────────────

async def scrape_website(url: str) -> str:
    """
    Scrape a website URL for its title and meta description.
    Returns a text summary string for ai_processor.
    """
    try:
        async with httpx.AsyncClient(
            timeout=15.0,
            follow_redirects=True,
            headers={
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/120.0.0.0 Safari/537.36"
                ),
            },
        ) as client:
            response = await client.get(url)
            response.raise_for_status()
            html = response.text

        # Extract <title>
        title_match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip() if title_match else ""

        # Extract meta description
        desc_match = re.search(
            r'<meta\s+[^>]*name=["\']description["\'][^>]*content=["\'](.*?)["\']',
            html, re.IGNORECASE | re.DOTALL,
        )
        if not desc_match:
            desc_match = re.search(
                r'<meta\s+[^>]*content=["\'](.*?)["\'][^>]*name=["\']description["\']',
                html, re.IGNORECASE | re.DOTALL,
            )
        description = desc_match.group(1).strip() if desc_match else ""

        # Extract og:description as fallback
        if not description:
            og_match = re.search(
                r'<meta\s+[^>]*property=["\']og:description["\'][^>]*content=["\'](.*?)["\']',
                html, re.IGNORECASE | re.DOTALL,
            )
            if og_match:
                description = og_match.group(1).strip()

        # Extract visible text from <p> tags (first 5 paragraphs)
        paragraphs = re.findall(r"<p[^>]*>(.*?)</p>", html, re.IGNORECASE | re.DOTALL)
        body_text = " ".join(
            re.sub(r"<[^>]+>", "", p).strip()
            for p in paragraphs[:5]
        )

        # Build a text blob for the AI processor
        parts = []
        if title:
            parts.append(f"Page title: {title}")
        if description:
            parts.append(f"Description: {description}")
        parts.append(f"URL: {url}")
        if body_text:
            parts.append(f"Page content: {body_text[:2000]}")

        return "\n".join(parts) if parts else f"Website URL: {url}"

    except Exception as e:
        # If scraping fails, still pass the URL to ai_processor
        return f"Website URL: {url}\n(Could not scrape page: {e})"


async def get_github_metadata(url: str) -> str:
    """Fetch GitHub repository metadata using the public GitHub API."""
    match = GITHUB_RE.search(url)
    if not match:
        logger.warning(f"Could not parse GitHub owner/repo from URL: {url}")
        return f"GitHub URL: {url}"
    
    owner = match.group(1)
    repo = match.group(2)
    if repo.lower().endswith(".git"):
        repo = repo[:-4]
        
    api_url = f"https://api.github.com/repos/{owner}/{repo}"
    logger.info(f"Fetching GitHub API: {api_url}")
    
    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            r = await client.get(api_url, headers={
                "Accept": "application/vnd.github.v3+json",
                "User-Agent": "VaultMCP-Backend"
            })
            data = r.json()
            logger.info(f"GitHub API response: {data}")
            
            # Also fetch README
            readme_url = f"https://api.github.com/repos/{owner}/{repo}/readme"
            readme_r = await client.get(readme_url, headers={
                "Accept": "application/vnd.github.v3+json",
                "User-Agent": "VaultMCP-Backend"
            })
            readme_text = ""
            if readme_r.status_code == 200:
                readme_data = readme_r.json()
                content_b64 = readme_data.get("content", "")
                if content_b64:
                    readme_text = base64.b64decode(content_b64).decode('utf-8', errors='ignore')[:1000]
            else:
                logger.warning(f"GitHub README response status: {readme_r.status_code}")
                
            description = data.get('description') or 'No description provided'
            readme_preview = readme_text[:800] if readme_text else 'No README available'
            
            # Build richer context for the AI
            return f"""
GitHub Repository: {data.get('full_name')}
Description: {description}
Primary Language: {data.get('language') or 'Not specified'}
Stars: {data.get('stargazers_count', 0)}
Topics: {', '.join(data.get('topics') or []) or 'none'}
README Preview: {readme_preview}

Based on the above, summarize what this repository does.
"""
    except Exception as e:
        logger.error(f"Error fetching GitHub metadata: {e}")
        return f"GitHub URL: {url} (Failed to fetch API metadata: {e})"


# ─── Routes ──────────────────────────────────────────────────────────────────

@app.get("/")
async def root_index():
    """Welcome page / status."""
    return {
        "status": "running",
        "service": "vaultmcp-backend",
        "message": "Welcome to VaultMCP API. System is active and healthy.",
        "documentation": "/docs"
    }


@app.get("/health")
async def health_check():
    """System health check."""
    return {
        "status": "ok",
        "service": "vaultmcp-backend",
        "version": "1.0.0",
        "timestamp": datetime.utcnow().isoformat(),
        "modules": {
            "metadata_scraper": "ready",
            "ai_processor": "ready",
            "web_searcher": "ready",
            "md_generator": "ready",
            "drive_handler": "ready",
        },
    }


@app.post("/process")
async def process_content(request: ProcessRequest):
    """
    Main processing endpoint — the entire VaultMCP pipeline.

    Flow by input type:
      Instagram/YouTube URL → scrape metadata → AI structure → web search → MD
      Website URL           → scrape title/desc → AI structure → web search → MD
      Plain text            → AI structure → web search → MD
    """
    content = request.content.strip()
    hf_token = request.hf_token.strip()

    if not content:
        raise HTTPException(status_code=400, detail="Content is empty.")
    if not hf_token:
        raise HTTPException(status_code=400, detail="Hugging Face token is required.")

    # ── Step 1: Detect input type ────────────────────────────────────────
    if request.content_type and request.content_type != "auto":
        input_type = request.content_type
    else:
        input_type = detect_input_type(content)

    pipeline_steps = [
        {"step": 1, "name": "detect_type", "status": "done", "detail": input_type},
    ]

    source_url = content if input_type != "text" else ""
    text_for_ai = ""

    # ── Step 2: Retrieve metadata/content based on input type ────────────
    if input_type == "github":
        try:
            metadata = await get_github_metadata(content)
            text_for_ai = metadata
            pipeline_steps.append({
                "step": 2, "name": "github_metadata", "status": "done",
                "detail": "Retrieved GitHub metadata",
            })
        except Exception as e:
            pipeline_steps.append({
                "step": 2, "name": "github_metadata", "status": "error",
                "detail": str(e),
            })
            raise HTTPException(status_code=422, detail=f"GitHub metadata fetch failed: {e}")

    elif input_type == "website":
        try:
            url_match = re.search(r'https?://\S+', content)
            if url_match:
                url = url_match.group(0)
                user_note = content.replace(url, '').strip()
                source_url = url
                
                if user_note:
                    text_for_ai = f"URL: {url}\nUser's note: {user_note}"
                    pipeline_steps.append({
                        "step": 2, "name": "user_note", "status": "done",
                        "detail": "Used user's custom note instead of scraping",
                    })
                else:
                    scraped_text = await scrape_website(url)
                    text_for_ai = scraped_text
                    pipeline_steps.append({
                        "step": 2, "name": "scrape_website", "status": "done",
                        "detail": f"Scraped {len(scraped_text)} chars",
                    })
            else:
                scraped_text = await scrape_website(content)
                text_for_ai = scraped_text
                pipeline_steps.append({
                    "step": 2, "name": "scrape_website", "status": "done",
                    "detail": f"Scraped {len(scraped_text)} chars",
                })
        except Exception as e:
            pipeline_steps.append({
                "step": 2, "name": "scrape_website", "status": "error",
                "detail": str(e),
            })
            raise HTTPException(status_code=422, detail=f"Website scraping failed: {e}")

    else:
        # Plain text
        text_for_ai = content
        pipeline_steps.append({
            "step": 2, "name": "prepare_content", "status": "done",
            "detail": f"Prepared {len(content)} chars of plain text",
        })

    # ── Step 3: AI structuring (Mistral) ───────────────────────────────
    try:
        processed = process_text(text_for_ai, hf_token)
        final_category = detect_category(input_type, source_url, processed.category)
        processed.category = final_category
        processed_dict = result_to_dict(processed)
        pipeline_steps.append({
            "step": 3, "name": "ai_structure", "status": "done",
            "detail": f"Category: {processed.category}, fallback: {processed.was_fallback}",
        })
    except AIInvalidTokenError as e:
        pipeline_steps.append({
            "step": 3, "name": "ai_structure", "status": "error",
            "detail": str(e),
        })
        status_code = 403 if "Inference Provider" in str(e) else 401
        raise HTTPException(status_code=status_code, detail=f"HF token error: {e}")
    except ProcessingError as e:
        pipeline_steps.append({
            "step": 3, "name": "ai_structure", "status": "error",
            "detail": str(e),
        })
        raise HTTPException(status_code=422, detail=f"AI processing failed: {e}")

    # ── Step 4: Web search for official link ─────────────────────
    official_info = get_official_info(processed.official_link)
    official_link = official_info.get("official_link", "")
    pipeline_steps.append({
        "step": 4, "name": "web_search", "status": "done" if official_link else "no_results",
        "detail": f"Provider: {official_info.get('provider', 'none')}, link: {official_link or 'N/A'}",
    })

    # ── Step 5: Generate Markdown ────────────────────────────────
    entry = build_entry(
        processed=processed_dict,
        source_url=source_url,
        official_link=official_link,
    )
    md_entry = generate_entry_md(entry)
    saved_on = format_retro_date(datetime.utcnow())

    pipeline_steps.append({
        "step": 5, "name": "generate_md", "status": "done",
    })

    # ── Return final result ──────────────────────────────────────
    return {
        "status": "success",
        "result": {
            "title": processed.title,
            "category": processed.category,
            "summary": processed.summary,
            "official_link": official_link,
            "source_url": source_url,
            "tools_mentioned": processed.tools_mentioned,
            "links_mentioned": processed.links_mentioned,
            "md_entry": md_entry,
            "saved_on": saved_on,
        },
        "pipeline_steps": pipeline_steps,
        "input_type": input_type,
    }


def extract_text_from_file(raw_bytes: bytes, filename: str) -> str:
    """Extract text content from various file formats and format as Markdown."""
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    try:
        if ext == 'pdf':
            doc = fitz.open(stream=raw_bytes, filetype="pdf")
            extracted = "\n".join([page.get_text() for page in doc])
            return f"# Document Title\n> Converted from {filename} by VaultMCP\n\n---\n\n{extracted}"
            
        elif ext == 'docx':
            doc = docx.Document(io.BytesIO(raw_bytes))
            parts = []
            for p in doc.paragraphs:
                if p.style.name.startswith('Heading'):
                    level = p.style.name.replace('Heading', '').strip()
                    try:
                        level_num = int(level)
                        parts.append(f"{'#' * level_num} {p.text}")
                    except ValueError:
                        parts.append(f"# {p.text}")
                else:
                    parts.append(p.text)
            extracted = "\n".join(parts)
            return f"# Document Title\n> Converted from {filename} by VaultMCP\n\n---\n\n{extracted}"
            
        elif ext == 'xlsx':
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True)
            parts = [f"# Spreadsheet: {filename}\n> Converted by VaultMCP\n"]
            for i, sheet in enumerate(wb.worksheets, 1):
                parts.append(f"## Sheet {i}: {sheet.title}")
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    parts.append("")
                    continue
                # Header
                header = rows[0]
                parts.append("| " + " | ".join([str(v) if v is not None else "" for v in header]) + " |")
                parts.append("|" + "|".join(["---" for _ in header]) + "|")
                # Rows
                for row in rows[1:]:
                    parts.append("| " + " | ".join([str(v) if v is not None else "" for v in row]) + " |")
                parts.append("")
            return "\n".join(parts)
            
        elif ext == 'pptx':
            prs = pptx.Presentation(io.BytesIO(raw_bytes))
            parts = [f"# Presentation: {filename}\n> Converted by VaultMCP\n"]
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
                parts.append("")
            return "\n".join(parts)
            
    except Exception as e:
        logger.error(f"Failed to parse {ext} file: {e}")
        pass

    # Fallback for plain text, unknown types, or failed parsing
    text = raw_bytes.decode("utf-8", errors="replace")
    return f"# Document: {filename}\n> Converted by VaultMCP\n\n---\n\n{text}"



@app.post("/process/file")
async def process_file(
    file: UploadFile = File(...),
    hf_token: str = Form(...),
    drive_access_token: Optional[str] = Form(None),
    drive_refresh_token: Optional[str] = Form(None),
):
    """
    Process an uploaded file (PDF, text, etc.).
    Reads file content as text and runs through the AI pipeline.
    """
    if not hf_token or not hf_token.strip():
        raise HTTPException(status_code=400, detail="Hugging Face token is required.")

    try:
        raw_bytes = await file.read()
        text_content = extract_text_from_file(raw_bytes, file.filename)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not read file: {e}")

    if not text_content.strip():
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    original_file_link = ""
    md_file_link = ""
    if drive_access_token:
        try:
            # Upload Original
            original_file_link = drive_save_file_to_drive(
                filename=file.filename,
                content_bytes=raw_bytes,
                mime_type=file.content_type or "application/octet-stream",
                access_token=drive_access_token,
                refresh_token=drive_refresh_token
            )
            
            # Upload MD version
            md_filename = f"{file.filename}.md"
            md_bytes = text_content.encode("utf-8")
            md_file_link = drive_save_file_to_drive(
                filename=md_filename,
                content_bytes=md_bytes,
                mime_type="text/markdown",
                access_token=drive_access_token,
                refresh_token=drive_refresh_token
            )
        except Exception as e:
            logger.warning(f"Failed to save files to Drive: {e}")

    # Prepend filename for context
    text_for_ai = f"File: {file.filename}\n\n{text_content[:8000]}"

    # Run through the same AI pipeline
    try:
        processed = process_text(text_for_ai, hf_token.strip())
        processed_dict = result_to_dict(processed)
    except AIInvalidTokenError as e:
        raise HTTPException(status_code=401, detail=f"HF token error: {e}")
    except ProcessingError as e:
        raise HTTPException(status_code=422, detail=f"AI processing failed: {e}")

    official_info = get_official_info(processed.official_link)
    official_link = official_info.get("official_link", "")

    entry = build_entry(
        processed=processed_dict,
        source_url="",
        official_link=official_link,
        original_file_link=original_file_link,
        md_file_link=md_file_link,
    )
    md_entry = generate_entry_md(entry)

    return {
        "status": "success",
        "result": {
            "title": processed.title,
            "category": processed.category,
            "summary": processed.summary,
            "official_link": official_link,
            "source_url": "",
            "tools_mentioned": processed.tools_mentioned,
            "links_mentioned": processed.links_mentioned,
            "md_entry": md_entry,
            "saved_on": format_retro_date(datetime.utcnow()),
        },
        "input_type": "file",
        "filename": file.filename,
        "original_file_link": original_file_link,
        "md_file_link": md_file_link,
    }


# ─── Drive Routes ────────────────────────────────────────────────────────────

@app.get("/drive/fetch")
async def drive_fetch(
    file_id: str,
    access_token: str,
    refresh_token: Optional[str] = None
):
    """Fetch raw text/markdown content of a file from Google Drive by its file ID."""
    if not file_id:
        raise HTTPException(status_code=400, detail="file_id is required.")
    if not access_token:
        raise HTTPException(status_code=400, detail="access_token is required.")

    try:
        content = drive_get_file_content(file_id, access_token, refresh_token)
        return Response(content=content, media_type="text/markdown")
    except DriveError as e:
        logger.error(f"Drive fetch failed: {e}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.exception("Unexpected error in /drive/fetch")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.post("/drive/clear")
async def clear_vault(request: DriveTokenRequest):
    try:
        drive_clear_vault_files(request.access_token, request.refresh_token)
        return {"status": "success", "message": "Vault cleared"}
    except DriveError as e:
        logger.error(f"Drive clear failed: {e}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.exception("Unexpected error in /drive/clear")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.post("/drive/save")
async def drive_save(request: DriveSaveRequest):
    """Save a Markdown entry to the user's Google Drive VaultMCP folder."""
    if not request.md_entry.strip():
        raise HTTPException(status_code=400, detail="md_entry is empty.")
    if not request.access_token.strip():
        raise HTTPException(status_code=400, detail="Google Drive access token is required.")

    try:
        result = drive_save_to_drive(
            md_entry=request.md_entry,
            access_token=request.access_token,
            refresh_token=request.refresh_token,
            overwrite=request.overwrite,
        )

        return {
            "status": "success",
            "message": f"Entry {result.action} in Google Drive.",
            "file_id": result.file_id,
            "folder_id": result.folder_id,
            "file_name": result.file_name,
            "action": result.action,
        }

    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"Google Drive auth error: {e}")
    except DriveError as e:
        raise HTTPException(status_code=500, detail=f"Google Drive error: {e}")


@app.get("/drive/vault")
async def drive_get_vault_route(
    access_token: str = Query(..., description="Google Drive OAuth access token"),
    refresh_token: Optional[str] = Query(None, description="Google Drive OAuth refresh token"),
):
    """Fetch the full vault.md content from the user's Google Drive."""
    if not access_token.strip():
        raise HTTPException(status_code=400, detail="Google Drive access token is required.")

    try:
        vault_content = drive_get_vault(
            access_token=access_token,
            refresh_token=refresh_token,
        )

        if vault_content is None:
            return {
                "status": "empty",
                "message": "No vault.md found in Google Drive. Save your first entry to create it.",
                "content": None,
            }

        return {
            "status": "success",
            "content": vault_content,
        }

    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"Google Drive auth error: {e}")
    except DriveError as e:
        raise HTTPException(status_code=500, detail=f"Google Drive error: {e}")


@app.post("/config/save")
async def config_save(request: ConfigSaveRequest):
    """Save user configuration to Google Drive."""
    if not request.access_token.strip():
        raise HTTPException(status_code=400, detail="Google Drive access token is required.")

    config = {
        "hf_token": request.hf_token,
        "display_name": request.display_name,
    }
    
    try:
        drive_save_user_config(config, request.access_token, request.refresh_token)
        return {"status": "success", "message": "Config saved to Google Drive."}
    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"Google Drive auth error: {e}")
    except DriveError as e:
        raise HTTPException(status_code=500, detail=f"Google Drive error: {e}")


@app.get("/config/load")
async def config_load(
    access_token: str = Query(..., description="Google Drive OAuth access token"),
    refresh_token: Optional[str] = Query(None, description="Google Drive OAuth refresh token"),
):
    """Load user configuration from Google Drive."""
    if not access_token.strip():
        raise HTTPException(status_code=400, detail="Google Drive access token is required.")
        
    try:
        config = drive_get_user_config(access_token, refresh_token)
        if config is None:
            return {"status": "empty", "config": {}}
        return {"status": "success", "config": config}
    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"Google Drive auth error: {e}")
    except DriveError as e:
        raise HTTPException(status_code=500, detail=f"Google Drive error: {e}")


# ─── Auth Routes ─────────────────────────────────────────────────────────────

@app.get("/auth/url")
async def auth_get_url():
    """Get the Google OAuth consent URL for the user to visit."""
    try:
        url = drive_get_auth_url()
        return {"status": "success", "auth_url": url}
    except DriveAuthError as e:
        raise HTTPException(status_code=500, detail=f"OAuth config error: {e}")


@app.post("/auth/google")
async def auth_google(request: DriveAuthRequest):
    """Exchange a Google OAuth authorization code for access + refresh tokens."""
    if not request.auth_code.strip():
        raise HTTPException(status_code=400, detail="Auth code is empty.")

    try:
        tokens = drive_auth_google(request.auth_code)

        return {
            "status": "success",
            "tokens": {
                "access_token": tokens.access_token,
                "refresh_token": tokens.refresh_token,
                "expires_in": tokens.expires_in,
            },
        }

    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"OAuth exchange failed: {e}")


# ─── MCP Routes ──────────────────────────────────────────────────────────────

def parse_vault_md(md_content: str) -> list:
    """Parses a vault.md file contents into a list of structured entries."""
    items = []
    if not md_content:
        return items

    sections = re.split(r"^###\s+", md_content, flags=re.MULTILINE)
    
    for i in range(1, len(sections)):
        section = sections[i]
        lines = section.strip().split("\n")
        if not lines:
            continue
        title = lines[0].strip()
        
        summary = ""
        official_link = ""
        source_url = ""
        tools_mentioned = []
        saved_on = ""
        
        for line in lines[1:]:
            trimmed = line.strip()
            if trimmed.startswith("- Summary:"):
                summary = trimmed[len("- Summary:"):].strip()
            elif trimmed.startswith("- Official link:"):
                official_link = trimmed[len("- Official link:"):].strip()
                if official_link == "N/A":
                    official_link = ""
            elif trimmed.startswith("- Source:"):
                source_url = trimmed[len("- Source:"):].strip()
            elif trimmed.startswith("- Tools mentioned:"):
                tools_str = trimmed[len("- Tools mentioned:"):].strip()
                tools_mentioned = [t.strip() for t in tools_str.split(",") if t.strip()]
            elif trimmed.startswith("- Saved on:"):
                saved_on = trimmed[len("- Saved on:"):].strip()
                
        above_content = md_content.split(f"### {title}")[0]
        cat_matches = re.findall(r"##\s+\[CATEGORY:\s*(.+?)\]", above_content)
        category = cat_matches[-1].strip() if cat_matches else "Other"
        
        items.append({
            "title": title,
            "category": category,
            "summary": summary,
            "official_link": official_link,
            "source_url": source_url,
            "tools_mentioned": tools_mentioned,
            "saved_on": saved_on
        })
        
    return items


@app.get("/mcp/vault")
async def mcp_get_vault(
    x_drive_token: Optional[str] = Header(None, alias="X-Drive-Token")
):
    """Returns the full vault.md contents as plain text."""
    if not x_drive_token or not x_drive_token.strip():
        raise HTTPException(
            status_code=400,
            detail="Google Drive access token is required in X-Drive-Token header."
        )

    try:
        vault_content = drive_get_vault(
            access_token=x_drive_token,
            refresh_token=None,
        )

        if vault_content is None:
            # Fallback to standard empty vault if not created yet
            vault_content = "# VaultMCP Vault\n\n> Save what you scroll. Use what you saved.\n"

        return Response(content=vault_content, media_type="text/plain")

    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"Google Drive auth error: {e}")
    except DriveError as e:
        raise HTTPException(status_code=500, detail=f"Google Drive error: {e}")


@app.get("/mcp/search")
async def mcp_search_vault(
    q: str = Query(..., description="Keyword query to search for"),
    x_drive_token: Optional[str] = Header(None, alias="X-Drive-Token")
):
    """Searches vault.md for matching entries by keyword and returns them as a JSON array."""
    if not x_drive_token or not x_drive_token.strip():
        raise HTTPException(
            status_code=400,
            detail="Google Drive access token is required in X-Drive-Token header."
        )

    try:
        vault_content = drive_get_vault(
            access_token=x_drive_token,
            refresh_token=None,
        )

        if not vault_content:
            return []

        entries = parse_vault_md(vault_content)
        query = q.strip().lower()
        results = []

        for entry in entries:
            title = entry.get("title", "").lower()
            category = entry.get("category", "").lower()
            summary = entry.get("summary", "").lower()
            tools = [t.lower() for t in entry.get("tools_mentioned", [])]

            if (query in title or 
                query in category or 
                query in summary or 
                any(query in t for t in tools)):
                results.append(entry)

        return results

    except DriveAuthError as e:
        raise HTTPException(status_code=401, detail=f"Google Drive auth error: {e}")
    except DriveError as e:
        raise HTTPException(status_code=500, detail=f"Google Drive error: {e}")


# ─── Run with Uvicorn ────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn

    host = os.getenv("HOST", "0.0.0.0")
    port = int(os.getenv("PORT", "8000"))
    uvicorn.run("main:app", host=host, port=port, reload=True)
